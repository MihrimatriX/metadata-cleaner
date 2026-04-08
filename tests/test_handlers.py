from pathlib import Path

import openpyxl
from pypdf import PdfReader, PdfWriter
from pptx import Presentation

from file_handlers.pdf_handler import remove_pdf_metadata
from file_handlers.pptx_handler import remove_pptx_metadata
from file_handlers.xlsx_handler import remove_xlsx_metadata
from remover import _skip_because_inside_output_only


def test_pdf_metadata_cleared(tmp_path: Path) -> None:
    src = tmp_path / "in.pdf"
    out = tmp_path / "out.pdf"
    w = PdfWriter()
    w.add_blank_page(width=72, height=72)
    w.add_metadata({"/Author": "Secret Author", "/Title": "Secret Title"})
    with open(src, "wb") as f:
        w.write(f)

    assert remove_pdf_metadata(str(src), str(out))

    reader = PdfReader(str(out))
    meta = reader.metadata or {}
    dumped = " ".join(str(v) for v in meta.values() if v is not None)
    assert "Secret" not in dumped


def test_xlsx_props_cleared(tmp_path: Path) -> None:
    src = tmp_path / "in.xlsx"
    out = tmp_path / "out.xlsx"
    wb = openpyxl.Workbook()
    wb.properties.creator = "Secret"
    wb.properties.title = "T"
    wb.save(src)

    assert remove_xlsx_metadata(str(src), str(out))

    wb2 = openpyxl.load_workbook(out)
    assert not (wb2.properties.creator or "").strip()
    assert not (wb2.properties.title or "").strip()


def test_pptx_props_cleared(tmp_path: Path) -> None:
    src = tmp_path / "in.pptx"
    out = tmp_path / "out.pptx"
    prs = Presentation()
    prs.slide_width = int(12 * 360000)
    prs.slide_height = int(7.5 * 360000)
    layouts = prs.slide_layouts
    layout = layouts[len(layouts) - 1]
    prs.slides.add_slide(layout)
    prs.core_properties.author = "Secret"
    prs.core_properties.title = "T"
    prs.save(src)

    assert remove_pptx_metadata(str(src), str(out))

    prs2 = Presentation(out)
    assert not (prs2.core_properties.author or "").strip()
    assert not (prs2.core_properties.title or "").strip()


def test_skip_output_subtree_only_when_distinct(tmp_path: Path) -> None:
    root = tmp_path / "project"
    cleaned = root / "cleaned"
    cleaned.mkdir(parents=True)
    inner = cleaned / "old.xlsx"
    inner.write_bytes(b"dummy")

    assert _skip_because_inside_output_only(str(inner), str(root), str(cleaned))
    assert not _skip_because_inside_output_only(str(root / "a.xlsx"), str(root), str(root))
